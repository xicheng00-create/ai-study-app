"""资料解析：PDF/PPTX/DOCX/MD/TXT → 纯文本 → 分块（RAG 降维替代）。

不引入向量嵌入；文本分块写 SQLite chunks，供 keyword 检索。
"""
import io

CHUNK_SIZE = 500
CHUNK_OVERLAP = 80

# 支持的文件类型 → 解析函数名
SUPPORTED = {"pdf", "pptx", "docx", "md", "txt", "markdown"}


def extract_text(filename: str, blob: bytes) -> str:
    """按扩展名分发解析，返回纯文本（失败抛异常，由调用方兜底）。"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in ("md", "markdown", "txt"):
        return _decode(blob)
    if ext == "pdf":
        return _extract_pdf(blob)
    if ext == "pptx":
        return _extract_pptx(blob)
    if ext == "docx":
        return _extract_docx(blob)
    raise ValueError(f"不支持的文件类型: {ext}")


def _decode(blob: bytes) -> str:
    for enc in ("utf-8", "gb18030", "latin-1"):
        try:
            return blob.decode(enc)
        except UnicodeDecodeError:
            continue
    return blob.decode("utf-8", errors="ignore")


def _extract_pdf(blob: bytes) -> str:
    import pdfplumber

    parts = []
    with pdfplumber.open(io.BytesIO(blob)) as pdf:
        for page in pdf.pages:
            txt = page.extract_text() or ""
            if txt:
                parts.append(txt)
    return "\n".join(parts)


def _extract_pptx(blob: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(blob))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_docx(blob: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(blob))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(parts)


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """按近似字符数滑窗分块（按段落优先，落到近似 size）。"""
    text = (text or "").strip()
    if not text:
        return []
    # 以换行/句号切粗块，避免从句子中间硬切
    paragraphs = [p.strip() for p in text.replace("\r", "\n").split("\n")]
    flat = "\n".join(p for p in paragraphs if p)
    if not flat:
        return []
    chunks = []
    start = 0
    n = len(flat)
    while start < n:
        end = min(start + size, n)
        if end < n:
            # 尽量退到最近的换行，避免截断句子
            nl = flat.rfind("\n", start, end)
            if nl > start + size // 2:
                end = nl + 1
        chunks.append(flat[start:end].strip())
        if end >= n:
            break
        start = max(start + size - overlap, end - overlap)
    return [c for c in chunks if c]


def chunk_text_list(text: str) -> list[dict]:
    """分块并带 chunk_idx（入库用）。"""
    return [{"chunk_idx": i, "text": c} for i, c in enumerate(chunk_text(text))]
